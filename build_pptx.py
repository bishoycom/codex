from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


BLUE = RGBColor(33, 150, 243)
NAVY = RGBColor(15, 32, 67)
LIGHT_BG = RGBColor(245, 249, 255)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(90, 98, 108)
ACCENT = RGBColor(0, 188, 212)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(25)
    p.font.color.rgb = WHITE

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.5))
        p2 = st.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = NAVY


def add_bullets(slide, items, left=0.8, top=1.5, width=6.2, height=4.8, font_size=22):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BLUE

    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(11)


def add_note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# 1 Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
hero.fill.solid(); hero.fill.fore_color.rgb = NAVY; hero.line.fill.background()
stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.8), Inches(13.333), Inches(1.7))
stripe.fill.solid(); stripe.fill.fore_color.rgb = BLUE; stripe.line.fill.background()

t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.8), Inches(2))
p = t.text_frame.paragraphs[0]; p.text = "Mastering SAT Math"; p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = t.text_frame.add_paragraph(); p2.text = "Strategies, Concepts, and Practice for the Digital SAT"; p2.font.size = Pt(24); p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER

for x in [1.2, 3.2, 5.2, 7.2, 9.2, 11.2]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(0.8), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
add_note(slide, "Welcome students and parents. Set the tone: SAT Math success comes from strategy + repetition.")

# 2 What is DSAT
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide, "What is the Digital SAT?")
add_bullets(slide, [
    "College admissions test accepted by most U.S. universities",
    "Fully digital format taken on a laptop/tablet",
    "Adaptive: Module 2 difficulty depends on Module 1 performance",
    "Math section: 2 modules, 35 min each, 44 questions total",
    "Built-in Desmos graphing calculator for all math questions",
], font_size=20)

# chart timing
cd = CategoryChartData()
cd.categories = ["Module 1", "Module 2"]
cd.add_series("Minutes", (35, 35))
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.3), Inches(1.7), Inches(5.4), Inches(3.3), cd).chart
chart.has_legend = False
chart.value_axis.maximum_scale = 40
add_note(slide, "Emphasize adaptive structure and equal module timing.")

# 3 topics
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide, "SAT Math Topics Overview")
cards = [
    ("📘 Algebra", "Linear equations, systems, inequalities"),
    ("📗 Advanced Math", "Quadratics, exponents, nonlinear functions"),
    ("📊 Problem Solving & Data", "Ratios, percentages, statistics"),
    ("📐 Geometry & Trig", "Angles, circles, area, right triangles"),
]
positions = [(0.8,1.5),(6.9,1.5),(0.8,4.2),(6.9,4.2)]
for (title, body),(x,y) in zip(cards,positions):
    box=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(5.6),Inches(2.2))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE; box.line.color.rgb=BLUE
    tf=box.text_frame
    p=tf.paragraphs[0]; p.text=title; p.font.bold=True; p.font.size=Pt(22); p.font.color.rgb=NAVY
    p2=tf.add_paragraph(); p2.text=body; p2.font.size=Pt(16); p2.font.color.rgb=GRAY

# 4 formulas
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"Most Important SAT Math Formulas")
formulae=[
"Linear: y = mx + b",
"Slope: m = (y₂ - y₁)/(x₂ - x₁)",
"Quadratic: x = (-b ± √(b²-4ac))/2a",
"Circle: C = 2πr,  A = πr²",
"Area: A=lw,  A=½bh,  A=πr²",
"Volume: V=lwh,  V=πr²h,  V=(4/3)πr³",
]
for i,f in enumerate(formulae):
    y=1.4+i*0.9
    b=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.9),Inches(y),Inches(11.8),Inches(0.7))
    b.fill.solid(); b.fill.fore_color.rgb=WHITE; b.line.color.rgb=BLUE
    p=b.text_frame.paragraphs[0]; p.text=f; p.font.size=Pt(20); p.font.color.rgb=NAVY

# 5 calc vs no
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"Calculator vs Mental Math Strategies")
add_bullets(slide,["Use mental math for simple arithmetic and percent changes","Use Desmos for intersections, roots, and visual checks","Estimate first to catch impossible answer choices","Avoid over-calculating: set up algebra before graphing","Common mistakes: wrong mode, typo inputs, rounding too early"],width=5.9,font_size=18)
add_bullets(slide,["No-calculator style habits still matter:","• Fraction/decimal fluency","• Factoring patterns","• Power rules","• Unit conversion speed"],left=7.0,width=5.4,font_size=18)

# 6 time management
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"SAT Math Time Management")
add_bullets(slide,["Pacing target: ~1.5 minutes/question","First pass: easy wins in first 12–15 minutes","Second pass: medium questions with setup time","Final pass: hard questions + informed guesses","Flag questions with long algebra or dense wording"],font_size=19)
cd=CategoryChartData(); cd.categories=["Easy","Medium","Hard"]; cd.add_series("Recommended Minutes",(15,12,8))
ch=slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,Inches(7.2),Inches(1.8),Inches(5.6),Inches(3.5),cd).chart
ch.has_legend=False

# 7 tricks
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"Common SAT Math Tricks")
add_bullets(slide,["Plug in numbers for variable-heavy word problems","Backsolve from answer choices (start with middle choice)","Estimate to eliminate outliers quickly","Use answer choice patterns (equivalent forms)","Example: If 3x+2=20, choices 4,5,6,7 → x=6"],font_size=19)

# 8 practice
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"SAT Practice Questions (with Solutions)")
qtext=[
"1) If 2x+5=19, x=?  A)5 B)6 C)7 D)8  → x=7",
"2) y=3x-4; find y when x=6. A)10 B)12 C)14 D)16 → y=14",
"3) x²-9x+20=0. A)4 or 5 B)2 or 10 C)-4 or -5 D)1 or 20 → 4,5",
"4) Circle radius 3: area=? A)6π B)9π C)12π D)18π → 9π",
"5) Data set 2,4,4,6,8 median=? A)4 B)4.8 C)5 D)6 → 4",
]
add_bullets(slide,qtext,left=0.6,top=1.3,width=12.1,height=5.8,font_size=16)
add_note(slide,"Work each question live. Ask students to justify each elimination step.")

# 9 desmos guide
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"Digital SAT Desmos Calculator Guide")
add_bullets(slide,["Graph equations quickly: y= forms and inequalities","Find intersections by clicking crossing points","Use tables to test values and patterns","Sliders visualize parameter changes","Check domain/range to avoid extraneous roots"],left=0.7,width=6.0,font_size=18)
mock=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(7.0),Inches(1.4),Inches(5.8),Inches(4.8))
mock.fill.solid(); mock.fill.fore_color.rgb=WHITE; mock.line.color.rgb=NAVY
p=mock.text_frame.paragraphs[0]; p.text="Desmos Mockup\n\ny=x^2-4\ny=2x\n\nIntersection: ( -1.24, -2.48 ) and (3.24, 6.48 )"; p.font.size=Pt(16); p.font.color.rgb=NAVY

#10 mistakes
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"Common Student Mistakes")
add_bullets(slide,["Misreading what the question asks (value vs expression)","Sign errors with negatives and distribution","Unit mistakes (inches vs feet, percent vs decimal)","Spending too long on one difficult problem","Skipping final answer reasonableness check"],font_size=19)

#11 study plan
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"SAT Study Plan")
add_bullets(slide,["1-Month Plan: 5 days/week, 60–90 min/day, topic sprints","3-Month Plan: concept phase → mixed sets → full tests","Weekly: 2 concept days, 2 drill days, 1 timed module day","Daily: 10 warm-up Qs + 20 targeted Qs + review mistakes","Track errors by category in an error log"],font_size=18)

#12 resources
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"Best SAT Resources")
add_bullets(slide,["Khan Academy Official SAT Practice","College Board Bluebook app + question bank","Official full-length practice tests","YouTube: Scalar Learning, John Jung, SuperTutorTV","Apps/Sites: UWorld, Quizlet SAT decks, Desmos tutorials"],font_size=18)

#13 motivation
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide,NAVY)
quote=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.1),Inches(1.5),Inches(11.1),Inches(4.5))
quote.fill.solid(); quote.fill.fore_color.rgb=WHITE; quote.line.color.rgb=BLUE
p=quote.text_frame.paragraphs[0]; p.text='“Success is the sum of small efforts, repeated day in and day out.”'; p.font.size=Pt(34); p.font.bold=True; p.font.color.rgb=NAVY; p.alignment=PP_ALIGN.CENTER
p2=quote.text_frame.add_paragraph(); p2.text='— Robert Collier'; p2.font.size=Pt(20); p2.font.color.rgb=GRAY; p2.alignment=PP_ALIGN.CENTER

#14 final
slide=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_header(slide,"Thank You")
end=slide.shapes.add_textbox(Inches(3.3),Inches(2.2),Inches(6.8),Inches(2.2))
p=end.text_frame.paragraphs[0]; p.text="Questions?"; p.font.size=Pt(54); p.font.bold=True; p.font.color.rgb=NAVY; p.alignment=PP_ALIGN.CENTER
p2=end.text_frame.add_paragraph(); p2.text="You’ve got this — practice smart and stay consistent."; p2.font.size=Pt(20); p2.font.color.rgb=BLUE; p2.alignment=PP_ALIGN.CENTER

prs.save("Mastering_SAT_Math_Digital_SAT.pptx")
print("Created Mastering_SAT_Math_Digital_SAT.pptx")
